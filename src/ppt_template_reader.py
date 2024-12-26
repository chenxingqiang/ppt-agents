from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import re


class PPTTemplateReader:
    def __init__(self, template_path):
        self.prs = Presentation(template_path)
        self.template_slides = self._categorize_slides()

    def _categorize_slides(self):
        categorized_slides = {
            "title": [],
            "toc": [],
            "content": [],
            "section": [],
            "closing": [],
        }
        for slide_layout in self.prs.slide_master.slide_layouts:
            slide_layout_name = slide_layout.name.lower()
            print(f"Slide layout name: {slide_layout_name}")
            if "title" in slide_layout_name:
                categorized_slides["title"].append(slide_layout)
            elif "toc" in slide_layout_name or "table of contents" in slide_layout_name:
                categorized_slides["toc"].append(slide_layout)
            elif "content" in slide_layout_name:
                categorized_slides["content"].append(slide_layout)
            elif "section" in slide_layout_name:
                categorized_slides["section"].append(slide_layout)
            elif "closing" in slide_layout_name:
                categorized_slides["closing"].append(slide_layout)
        print(f"Categorized slides: {categorized_slides}")
        return categorized_slides

    def create_new_slide(self, category, title_text, content_text):
        if category not in self.template_slides:
            raise ValueError(f"No slides found for category '{category}'")

        slide_layout_list = self.template_slides[category]
        if not slide_layout_list:
            raise ValueError(f"No slides available for category '{category}'")

        slide_layout = slide_layout_list[0]
        slide = self.prs.slides.add_slide(slide_layout)

        # Set the title
        if slide.shapes.title:
            slide.shapes.title.text = title_text

        # Find the content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type in [1, 2, 3, 13, 14, 15]:  # Common content placeholder types
                content_placeholder = shape
                break

        if content_placeholder and content_text:
            self._set_content(content_placeholder, content_text)
        elif content_text:
            # If no placeholder found, create a new textbox
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            self._set_content(textbox, content_text)

        return slide

    def _set_content(self, shape, content_text):
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        # Split content into paragraphs
        paragraphs = content_text.split('\n')
        
        for para_text in paragraphs:
            if not para_text.strip():
                continue
                
            p = text_frame.add_paragraph()
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            
            # Handle bullet points
            if para_text.startswith('•'):
                p.level = 1
                para_text = para_text[1:].strip()
            
            # Add the text with proper formatting
            run = p.add_run()
            run.text = para_text
            
            # Set font properties
            font = run.font
            font.name = 'Arial'
            font.size = Pt(18)
            
            # If text contains a URL, make it blue
            if 'http://' in para_text or 'https://' in para_text:
                font.color.rgb = RGBColor(0, 0, 255)

    def save_presentation(self, output_path):
        try:
            self.prs.save(output_path)
            return True
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return False
