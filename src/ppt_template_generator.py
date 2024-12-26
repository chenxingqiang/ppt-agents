from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PPTTemplateGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.slide_width = Inches(13.33)
        self.slide_height = Inches(7.5)
        
    def set_slide_size(self, width_inches=13.33, height_inches=7.5):
        """Set slide size in inches (default is 16:9)"""
        self.slide_width = Inches(width_inches)
        self.slide_height = Inches(height_inches)
        self.prs.slide_width = self.slide_width
        self.prs.slide_height = self.slide_height

    def _add_title_layout(self):
        """添加标题页布局"""
        # 使用默认布局
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        
        # 添加标题占位符
        title = slide.shapes.add_textbox(
            Inches(1),  # left
            Inches(2),  # top
            Inches(11.33),  # width
            Inches(2)  # height
        )
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 添加副标题占位符
        subtitle = slide.shapes.add_textbox(
            Inches(2),  # left
            Inches(4.5),  # top
            Inches(9.33),  # width
            Inches(1)  # height
        )
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_section_layout(self):
        """添加章节页布局"""
        # 使用默认布局
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        
        # 添加大标题占位符
        title = slide.shapes.add_textbox(
            Inches(1),
            Inches(3),
            Inches(11.33),
            Inches(1.5)
        )
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_content_layout(self):
        """添加内容页布局"""
        # 使用默认布局
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        
        # 添加标题占位符
        title = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.5),
            Inches(12.33),
            Inches(1)
        )
        
        # 添加内容占位符
        content = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(12.33),
            Inches(5.5)
        )

    def _add_two_column_layout(self):
        """添加双栏布局"""
        # 使用默认布局
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        
        # 添加标题占位符
        title = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.5),
            Inches(12.33),
            Inches(1)
        )
        
        # 添加左栏占位符
        left_content = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(6),
            Inches(5.5)
        )
        
        # 添加右栏占位符
        right_content = slide.shapes.add_textbox(
            Inches(6.83),
            Inches(1.5),
            Inches(6),
            Inches(5.5)
        )

    def _add_comparison_layout(self):
        """添加对比布局"""
        # 使用默认布局
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        
        # 添加标题占位符
        title = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.5),
            Inches(12.33),
            Inches(1)
        )
        
        # 左侧内容
        left_title = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(6),
            Inches(1)
        )
        left_content = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(2.5),
            Inches(6),
            Inches(4.5)
        )
        
        # 右侧内容
        right_title = slide.shapes.add_textbox(
            Inches(6.83),
            Inches(1.5),
            Inches(6),
            Inches(1)
        )
        right_content = slide.shapes.add_textbox(
            Inches(6.83),
            Inches(2.5),
            Inches(6),
            Inches(4.5)
        )

    def _add_image_with_text_layout(self):
        """添加图文混排布局"""
        # 使用默认布局
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        
        # 添加标题占位符
        title = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.5),
            Inches(12.33),
            Inches(1)
        )
        
        # 添加图片占位符
        image = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(6),
            Inches(5.5)
        )
        
        # 添加文本占位符
        text = slide.shapes.add_textbox(
            Inches(6.83),
            Inches(1.5),
            Inches(6),
            Inches(5.5)
        )

    def generate_template(self, theme="modern", output_path="template.pptx"):
        """生成PPT模板"""
        # 添加各种布局
        self._add_title_layout()
        self._add_section_layout()
        self._add_content_layout()
        self._add_two_column_layout()
        self._add_comparison_layout()
        self._add_image_with_text_layout()
        
        # 应用主题样式
        self._apply_theme(theme)
        
        # 保存模板
        self.prs.save(output_path)
        return output_path
    
    def _apply_theme(self, theme):
        """应用主题样式"""
        if theme == "modern":
            self._apply_modern_theme()
        elif theme == "classic":
            self._apply_classic_theme()
        elif theme == "minimal":
            self._apply_minimal_theme()
        elif theme == "tech":
            self._apply_tech_theme()
        elif theme == "nature":
            self._apply_nature_theme()
        elif theme == "elegant":
            self._apply_elegant_theme()
        elif theme == "vibrant":
            self._apply_vibrant_theme()
        elif theme == "corporate":
            self._apply_corporate_theme()
        elif theme == "creative":
            self._apply_creative_theme()
        elif theme == "dark":
            self._apply_dark_theme()
        else:
            self._apply_default_theme()
    
    def _apply_modern_theme(self):
        """现代主题 - 清新科技风"""
        theme_colors = [
            RGBColor(0, 114, 198),  # 主色：深蓝色
            RGBColor(0, 175, 240),  # 辅助色1：浅蓝色
            RGBColor(127, 186, 0),  # 辅助色2：青绿色
            RGBColor(255, 140, 0),  # 强调色1：橙色
            RGBColor(255, 67, 67),  # 强调色2：红色
        ]
        self._apply_theme_colors(theme_colors, "Arial", 40, 24)

    def _apply_tech_theme(self):
        """科技主题 - 未来感"""
        theme_colors = [
            RGBColor(28, 28, 35),   # 主色：深空灰
            RGBColor(75, 207, 250),  # 辅助色1：霓虹蓝
            RGBColor(52, 231, 228),  # 辅助色2：青色
            RGBColor(255, 94, 247),  # 强调色1：霓虹粉
            RGBColor(93, 248, 152),  # 强调色2：霓虹绿
        ]
        self._apply_theme_colors(theme_colors, "Helvetica", 42, 24)

    def _apply_nature_theme(self):
        """自然主题 - 生态环保"""
        theme_colors = [
            RGBColor(76, 175, 80),   # 主色：翠绿色
            RGBColor(139, 195, 74),  # 辅助色1：草绿色
            RGBColor(205, 220, 57),  # 辅助色2：柠檬绿
            RGBColor(255, 235, 59),  # 强调色1：向日葵黄
            RGBColor(255, 193, 7),   # 强调色2：琥珀色
        ]
        self._apply_theme_colors(theme_colors, "Georgia", 38, 22)

    def _apply_elegant_theme(self):
        """优雅主题 - 高端大气"""
        theme_colors = [
            RGBColor(63, 63, 63),    # 主色：深灰色
            RGBColor(158, 118, 80),  # 辅助色1：古铜色
            RGBColor(213, 196, 161), # 辅助色2：米色
            RGBColor(123, 63, 0),    # 强调色1：深棕色
            RGBColor(168, 168, 168), # 强调色2：银灰色
        ]
        self._apply_theme_colors(theme_colors, "Garamond", 36, 20)

    def _apply_vibrant_theme(self):
        """活力主题 - 充满朝气"""
        theme_colors = [
            RGBColor(255, 87, 51),   # 主色：珊瑚红
            RGBColor(255, 189, 74),  # 辅助色1：明黄色
            RGBColor(255, 136, 136), # 辅助色2：粉红色
            RGBColor(102, 187, 255), # 强调色1：天蓝色
            RGBColor(171, 140, 255), # 强调色2：淡紫色
        ]
        self._apply_theme_colors(theme_colors, "Verdana", 40, 24)

    def _apply_corporate_theme(self):
        """企业主题 - 专业稳重"""
        theme_colors = [
            RGBColor(25, 51, 92),    # 主色：海军蓝
            RGBColor(44, 62, 80),    # 辅助色1：深蓝灰
            RGBColor(52, 152, 219),  # 辅助色2：商务蓝
            RGBColor(155, 89, 182),  # 强调色1：优雅紫
            RGBColor(231, 76, 60),   # 强调色2：商务红
        ]
        self._apply_theme_colors(theme_colors, "Calibri", 38, 22)

    def _apply_creative_theme(self):
        """创意主题 - 艺术感"""
        theme_colors = [
            RGBColor(241, 90, 135),  # 主色：玫瑰粉
            RGBColor(255, 159, 26),  # 辅助色1：橘黄色
            RGBColor(88, 177, 159),  # 辅助色2：薄荷绿
            RGBColor(116, 95, 181),  # 强调色1：梦幻紫
            RGBColor(54, 215, 183),  # 强调色2：绿松石
        ]
        self._apply_theme_colors(theme_colors, "Century Gothic", 40, 24)

    def _apply_dark_theme(self):
        """暗黑主题 - 高端酷炫"""
        theme_colors = [
            RGBColor(18, 18, 18),    # 主色：纯黑
            RGBColor(45, 45, 45),    # 辅助色1：深灰
            RGBColor(83, 83, 83),    # 辅助色2：中灰
            RGBColor(0, 255, 255),   # 强调色1：青色
            RGBColor(255, 64, 129),  # 强调色2：粉红
        ]
        self._apply_theme_colors(theme_colors, "Consolas", 40, 24)

    def _apply_theme_colors(self, theme_colors, font_name, title_size, content_size):
        """统一的主题应用方法"""
        for slide in self.prs.slides:
            # 设置背景色（如果需要）
            if hasattr(slide, "background"):
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = theme_colors[1]
            
            # 应用到形状
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = font_name
                            if shape.name == "Title" or "title" in shape.name.lower():
                                run.font.size = Pt(title_size)
                                run.font.color.rgb = theme_colors[0]
                            else:
                                run.font.size = Pt(content_size)
                                run.font.color.rgb = theme_colors[2]
                                
                # 设置形状填充色（如果有）
                if hasattr(shape, "fill"):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = theme_colors[1] 