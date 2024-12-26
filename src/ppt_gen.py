from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import re
import os


def remove_non_alphanumeric(s):
    # 只保留中文、英文和数字字符
    filtered_string = re.sub(r"[^\u4e00-\u9fa5a-zA-Z0-9]", "", s)
    return filtered_string

from datetime import datetime

class DetailedPPTGenerator:
    def __init__(self, template_reader, markdown_reader, output_dir):
        self.template_reader = template_reader
        self.markdown_reader = markdown_reader
        self.output_dir = output_dir
        self.title = markdown_reader.get_title()
        self.subtitle = markdown_reader.get_subtitle()

    def _format_references(self, references):
        """Format references for slide notes."""
        if not references:
            return ""
        
        # Group references by section
        formatted_refs = []
        formatted_refs.append("参考文献：")
        formatted_refs.append("-" * 40)  # Add a separator line
        
        for ref in references:
            # Add some spacing for better readability
            formatted_refs.append("")  # Add blank line before each reference
            formatted_refs.append(ref)
        
        return "\n".join(formatted_refs)

    def _add_slide(self, category, title, content, notes=None, references=None):
        """Add a slide with content and references in notes."""
        slide = self.template_reader.create_new_slide(category, title, content)
        if slide.notes_slide:
            notes_text = []
            
            # Add content notes if any
            if notes:
                notes_text.extend(notes)
                notes_text.append("")  # Add a blank line after notes
            
            # Add references if any
            if references:
                notes_text.append(self._format_references(references))
            
            if notes_text:
                slide.notes_slide.notes_text_frame.text = "\n".join(notes_text)
        
        return slide

    def _format_table_for_slide(self, table_data):
        """Format table data into a string suitable for PPT."""
        lines = []
        
        # Add headers
        if table_data["headers"]:
            headers = " | ".join(table_data["headers"])
            lines.append(headers)
            lines.append("-" * len(headers))  # Add separator
        
        # Add rows
        for row in table_data["rows"]:
            lines.append(" | ".join(row))
        
        return "\n".join(lines)

    def generate_ppt(self):
        # Title Slide
        self._add_slide("title", self.title, self.subtitle)

        # Get all sections
        sections = self.markdown_reader.get_sections()
        
        # Generate Table of Contents
        toc_content = "\n".join(
            [f"{i+1}. {section['title']}"
             for i, section in enumerate(sections)]
        )
        self._add_slide("toc", "目录", toc_content)

        # Process each section
        for section in sections:
            # Add section title slide
            self._add_slide("section", section["title"], "", references=section.get("references", []))

            # Process content based on type
            for content_item in section["content"]:
                if content_item["type"] == "table":
                    # Format table for slide
                    table_content = self._format_table_for_slide(content_item["data"])
                    self._add_slide(
                        "content",
                        section["title"],
                        table_content,
                        notes=section.get("notes", []),
                        references=section.get("references", [])
                    )
                elif content_item["type"] == "list":
                    # Format list items
                    list_content = "\n".join(content_item["data"])
                    self._add_slide(
                        "content",
                        section["title"],
                        list_content,
                        notes=section.get("notes", []),
                        references=section.get("references", [])
                    )
                else:  # text content
                    self._add_slide(
                        "content",
                        section["title"],
                        content_item["data"],
                        notes=section.get("notes", []),
                        references=section.get("references", [])
                    )

        # Save the presentation
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)

        # Generate timestamp for versioning
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Create filename
        markdown_filename = os.path.basename(self.markdown_reader.markdown_path)
        markdown_name = os.path.splitext(markdown_filename)[0]
        sanitized_title = remove_non_alphanumeric(self.title)
        sanitized_markdown_name = remove_non_alphanumeric(markdown_name)
        filename = f"{timestamp}-{sanitized_markdown_name}-{sanitized_title}.pptx"
        
        output_path = os.path.join(self.output_dir, filename)
        self.template_reader.save_presentation(output_path)
        return output_path