from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from typing import Dict, List, Tuple
import re
from src.layout_analyzer import LayoutAnalyzer

class PPTGenerator:
    def __init__(self, template_path: str = None):
        self.template_path = template_path or "templates/modern_template.pptx"
        self.prs = Presentation(self.template_path)
        self.layout_analyzer = LayoutAnalyzer()
        self.slide_layouts = self._categorize_slide_layouts()
        
        # Define font sizes
        self.font_sizes = {
            1: Pt(40),  # # Heading 1
            2: Pt(32),  # ## Heading 2
            3: Pt(28),  # ### Heading 3
            4: Pt(24),  # #### Heading 4
            'body': Pt(18),  # Normal text
            'bullet_l1': Pt(18),
            'bullet_l2': Pt(16),
            'bullet_l3': Pt(14),
            'table_header': Pt(16),
            'table_cell': Pt(14),
        }
        
        # Define table styles
        self.table_styles = {
            'academic': {
                'header_bg': RGBColor(69, 69, 69),  # Dark gray
                'header_font': RGBColor(255, 255, 255),  # White
                'row_bg_even': RGBColor(242, 242, 242),  # Light gray
                'row_bg_odd': RGBColor(255, 255, 255),  # White
                'cell_font': RGBColor(0, 0, 0),  # Black
                'border_color': RGBColor(128, 128, 128),  # Gray
            }
        }

    def _categorize_slide_layouts(self) -> Dict:
        """Categorize available slide layouts by type."""
        layouts = {}
        for layout in self.prs.slide_layouts:
            name = layout.name.lower()
            if 'title' in name and 'only' in name:
                layouts['title'] = layout
            elif 'section' in name:
                layouts['section'] = layout
            elif 'two' in name and 'content' in name:
                layouts['two_column'] = layout
            elif 'comparison' in name:
                layouts['comparison'] = layout
            elif 'content' in name:
                layouts['content'] = layout
            elif 'blank' in name:
                layouts['blank'] = layout
        return layouts

    def _get_best_layout(self, layout_type: str):
        """Get the best matching slide layout for the content type."""
        if layout_type in self.slide_layouts:
            return self.slide_layouts[layout_type]
        return self.slide_layouts.get('content', self.prs.slide_layouts[1])  # default layout

    def _get_heading_level(self, line: str) -> int:
        """Get the heading level from a markdown line."""
        if line.startswith('#'):
            return len(line.split()[0].strip())
        return 0

    def _parse_table(self, table_text: str) -> Tuple[List[str], List[List[str]]]:
        """Parse markdown table into headers and rows."""
        lines = table_text.strip().split('\n')
        if len(lines) < 3:  # Need at least header, separator, and one row
            return [], []
            
        # Parse headers
        headers = [cell.strip() for cell in lines[0].strip('|').split('|')]
        
        # Skip separator line
        rows = []
        for line in lines[2:]:  # Skip header and separator
            if line.strip():
                cells = [cell.strip() for cell in line.strip('|').split('|')]
                rows.append(cells)
                
        return headers, rows

    def _calculate_content_size(self, content: str) -> tuple:
        """Calculate optimal size for content block based on content length and type."""
        lines = content.strip().split('\n')
        max_line_length = max(len(line) for line in lines)
        line_count = len(lines)
        
        # Base width calculation (assuming average character width)
        width = min(max_line_length * 0.1, 8.0)  # Max width of 8 inches
        
        # Base height calculation
        height = min(line_count * 0.3, 5.0)  # Max height of 5 inches
        
        return width, height

    def _calculate_table_size(self, headers: List[str], rows: List[List[str]]) -> tuple:
        """Calculate optimal table size based on content."""
        # Calculate column widths based on content
        col_widths = []
        for col_idx in range(len(headers)):
            col_content = [headers[col_idx]] + [row[col_idx] for row in rows if col_idx < len(row)]
            max_width = max(len(str(cell)) for cell in col_content)
            col_widths.append(max_width)
        
        # Calculate total width (with some padding)
        total_chars = sum(col_widths)
        table_width = min(total_chars * 0.15, 9.0)  # Max width of 9 inches
        
        # Calculate row height based on content
        max_lines_in_cell = 1
        for row in rows:
            for cell in row:
                lines_in_cell = len(str(cell).split('\n'))
                max_lines_in_cell = max(max_lines_in_cell, lines_in_cell)
        
        # Calculate total height
        row_height = 0.3  # Base height per row in inches
        total_rows = len(rows) + 1  # Include header
        table_height = min(total_rows * row_height * max_lines_in_cell, 5.0)  # Max height of 5 inches
        
        # Calculate individual column widths proportionally
        proportional_widths = [width/total_chars * table_width for width in col_widths]
        
        return table_width, table_height, proportional_widths

    def _add_table_to_slide(self, slide, table_text: str, position: tuple):
        """Add a formatted table to the slide."""
        headers, rows = self._parse_table(table_text)
        if not headers or not rows:
            return
            
        # Calculate optimal table dimensions
        table_width, table_height, col_widths = self._calculate_table_size(headers, rows)
        
        # Create table with calculated dimensions
        x, y, _, _ = position  # Ignore provided width/height, use calculated ones
        left = Inches(x * 10)
        top = Inches(y * 7.5)
        
        shape = slide.shapes.add_table(
            len(rows) + 1,  # Include header row
            len(headers),
            left,
            top,
            Inches(table_width),
            Inches(table_height)
        )
        table = shape.table
        
        # Set column widths
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)
        
        # Apply academic style
        style = self.table_styles['academic']
        
        # Format headers
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = style['header_bg']
            
            # Set vertical alignment to center
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.text = header
            
            run = paragraph.runs[0]
            run.font.size = self.font_sizes['table_header']
            run.font.bold = True
            run.font.color.rgb = style['header_font']
            
        # Format rows
        for row_idx, row_data in enumerate(rows, start=1):
            # Set row height based on content
            max_lines = max(len(str(cell).split('\n')) for cell in row_data)
            table.rows[row_idx].height = Inches(max_lines * 0.3)
            
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                
                # Alternate row background colors
                cell.fill.solid()
                cell.fill.fore_color.rgb = style['row_bg_even'] if row_idx % 2 == 0 else style['row_bg_odd']
                
                # Set vertical alignment to center
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.text = str(cell_text)
                
                run = paragraph.runs[0]
                run.font.size = self.font_sizes['table_cell']
                run.font.color.rgb = style['cell_font']
                
        # Apply borders
        for row_idx in range(len(rows) + 1):
            for col_idx in range(len(headers)):
                cell = table.cell(row_idx, col_idx)
                for border in ['top', 'right', 'bottom', 'left']:
                    line = getattr(cell.borders, border)
                    line.color.rgb = style['border_color']
                    line.width = Pt(1)

    def _apply_content_block(self, shape, content: str, position: tuple):
        """Apply content to a shape with specified position."""
        if not shape:
            return

        # Check if content is a table
        if content.strip().startswith('|') and content.strip().endswith('|'):
            self._add_table_to_slide(shape.part.slide, content, position)
            return
        
        # Calculate optimal content size
        width, height = self._calculate_content_size(content)
        
        # Apply position and size
        x, y, _, _ = position  # Ignore provided width/height, use calculated ones
        shape.left = Inches(x * 10)
        shape.top = Inches(y * 7.5)
        shape.width = Inches(width * 10)
        shape.height = Inches(height * 7.5)
        
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Use correct auto-size enumeration
        
        # Process normal content
        lines = content.strip().split('\n')
        current_level = 0
        
        for line in lines:
            if not line.strip():
                continue
                
            p = text_frame.add_paragraph()
            
            # Handle headings
            heading_level = self._get_heading_level(line)
            if heading_level > 0:
                p.text = line.lstrip('#').strip()
                p.font.size = self.font_sizes.get(heading_level, self.font_sizes['body'])
                p.font.bold = True
                continue
            
            # Handle bullet points and indentation
            stripped_line = line.lstrip()
            indent_level = len(line) - len(stripped_line)
            bullet_level = indent_level // 2
            
            if stripped_line.startswith('- ') or stripped_line.startswith('* '):
                p.text = stripped_line[2:].strip()
                p.level = min(bullet_level, 2)
                p.bullet = True
                p.font.size = self.font_sizes.get(f'bullet_l{p.level + 1}', self.font_sizes['body'])
            else:
                p.text = line.strip()
                p.font.size = self.font_sizes['body']

    def create_slide(self, markdown_content: str):
        """Create a single slide from markdown content."""
        # Split content into title and body
        lines = markdown_content.strip().split('\n')
        
        # Find title (## heading)
        title = None
        content_start = 0
        layout_type = 'content'  # default layout
        
        for i, line in enumerate(lines):
            if line.startswith('## '):
                title = line.strip('# ')
                content_start = i + 1
                break
            elif line.startswith('# '):
                title = line.strip('# ')
                content_start = i + 1
                layout_type = 'title'
                break
        
        if not title:
            return
            
        # Get content
        content = '\n'.join(lines[content_start:])
        
        # Get appropriate slide layout
        slide_layout = self._get_best_layout(layout_type)
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Apply title with appropriate font size
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
                title_shape.text_frame.paragraphs[0].font.size = self.font_sizes[1 if layout_type == 'title' else 2]
        
        # Apply content
        if hasattr(slide, 'placeholders'):
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Main content placeholder
                    self._apply_content_block(shape, content, (0.1, 0.3, 0.8, 0.6))
                    break

    def process_markdown_file(self, markdown_path: str, output_path: str):
        """Process a complete markdown file and generate PowerPoint presentation."""
        with open(markdown_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Split content into slides
        slides_content = content.split('---')
        
        # Process each slide
        for slide_content in slides_content:
            if slide_content.strip():
                self.create_slide(slide_content.strip())
        
        # Save presentation
        self.prs.save(output_path)

    def apply_theme(self, theme_name: str):
        """Apply a predefined theme to the presentation."""
        theme_settings = {
            'modern': {
                'background_color': 'FFFFFF',
                'title_font': 'Arial',
                'body_font': 'Arial',
                'accent_color': '4472C4'
            },
            'classic': {
                'background_color': 'F2F2F2',
                'title_font': 'Times New Roman',
                'body_font': 'Calibri',
                'accent_color': '8B0000'
            }
        }
        
        if theme_name in theme_settings:
            theme = theme_settings[theme_name]
            
            # Apply theme settings to slide master
            for layout in self.prs.slide_layouts:
                for shape in layout.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                if 'title' in layout.name.lower():
                                    run.font.name = theme['title_font']
                                else:
                                    run.font.name = theme['body_font'] 